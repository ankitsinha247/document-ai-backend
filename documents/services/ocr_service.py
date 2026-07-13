import os
import fitz
import pytesseract
import pandas as pd
from PIL import Image
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image, ImageOps


# -------------------------------------------------------
# PDF TEXT EXTRACTION
# -------------------------------------------------------

def extract_pdf_text(pdf_path):

    doc = fitz.open(pdf_path)

    full_text = ""

    custom_config = (
        "--oem 3 "
        "--psm 6 "
        "-l eng"
    )

    for page_num in range(len(doc)):

        page = doc.load_page(page_num)

        # ---------------------------------
        # Try native PDF text first
        # ---------------------------------

        text = page.get_text("text")

        if len(text.strip()) > 50:

            full_text += text + "\n"

            continue

        print(f"OCR Page : {page_num + 1}")

        # ---------------------------------
        # Convert PDF to High Resolution Image
        # ---------------------------------

        pix = page.get_pixmap(
            matrix=fitz.Matrix(4, 4)
        )

        image_path = f"/tmp/page_{page_num}.png"

        pix.save(image_path)

        image = Image.open(image_path)

        # ---------------------------------
        # OCR Pre-processing
        # ---------------------------------

        image = ImageOps.exif_transpose(image)

        image = ImageOps.grayscale(image)

        image = ImageOps.autocontrast(image)

        image = image.point(
            lambda p: 255 if p > 170 else 0
        )

        text = pytesseract.image_to_string(
            image,
            config=custom_config
        )

        full_text += text + "\n"

        if os.path.exists(image_path):
            os.remove(image_path)

    doc.close()

    return full_text.strip()
""" def extract_pdf_text(pdf_path):

    doc = fitz.open(pdf_path)

    full_text = ""

    for page_num in range(len(doc)):

        page = doc.load_page(page_num)

        text = page.get_text()

        if text.strip():

            full_text += text

        else:

            pix = page.get_pixmap(matrix=fitz.Matrix(2,2))

            image_path = f"/tmp/page_{page_num}.png"

            pix.save(image_path)

            image = Image.open(image_path)

            full_text += pytesseract.image_to_string(image)

    return full_text.strip() """
    


# -------------------------------------------------------
# IMAGE OCR
# -------------------------------------------------------

def extract_image_text(image_path):

    print("Reading Image :", image_path)

    image = Image.open(image_path)
    image = ImageOps.exif_transpose(image)
    image = image.convert("RGB")
    image = ImageOps.grayscale(image)
    image = ImageOps.autocontrast(image)
    image = image.point(
        lambda p: 255 if p > 170 else 0
    )

    text = pytesseract.image_to_string(
        image,
        config="--oem 3 --psm 6 -l eng"
    )

    print("OCR TEXT LENGTH :", len(text))

    return text.strip()
""" 
def extract_image_text(image_path):

    print("Reading Image:", image_path)

    image = Image.open(image_path)

    # Convert to grayscale
    image = ImageOps.grayscale(image)

    # Increase OCR accuracy
    image = image.point(lambda x: 0 if x < 140 else 255)

    text = pytesseract.image_to_string(
        image,
        config="--psm 6"
    )

    print("OCR TEXT LENGTH:", len(text))

    return text.strip()
 """

# -------------------------------------------------------
# WORD DOCUMENT
# -------------------------------------------------------

def extract_docx_text(file_path):

    doc = Document(file_path)

    text = ""

    for para in doc.paragraphs:

        text += para.text + "\n"

    return text


# -------------------------------------------------------
# EXCEL
# -------------------------------------------------------

def extract_excel_text(file_path):

    workbook = load_workbook(file_path)

    text = ""

    for sheet in workbook.worksheets:

        for row in sheet.iter_rows(values_only=True):

            text += " ".join(
                [str(cell) for cell in row if cell]
            )

            text += "\n"

    return text


# -------------------------------------------------------
# CSV
# -------------------------------------------------------

def extract_ppt_text(file_path):

    prs = Presentation(file_path)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text


# -------------------------------------------------------
# TXT
# -------------------------------------------------------

def extract_csv_text(file_path):

    df = pd.read_csv(file_path)

    return df.to_string()



def extract_txt_text(file_path):

    with open(file_path,"r",encoding="utf-8") as f:

        return f.read()
    

# -------------------------------------------------------
# MAIN UNIVERSAL FUNCTION
# -------------------------------------------------------
def extract_text(file_path, document_type):

    document_type = document_type.lower()
    extension = os.path.splitext(file_path)[1].lower()

    print("DOCUMENT TYPE :", document_type)
    print("FILE PATH :", file_path)

    image_types = {
        "image",
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".tif",
        ".tiff",
        "png",
        "jpg",
        "jpeg",
        "bmp",
        "tif",
        "tiff",
    }

    if document_type == "pdf":
        return extract_pdf_text(file_path)

    elif document_type in image_types or extension in image_types:
        return extract_image_text(file_path)

    elif document_type == "docx":
        return extract_docx_text(file_path)

    elif document_type == "xlsx":
        return extract_excel_text(file_path)

    elif document_type == "pptx":
        return extract_ppt_text(file_path)

    elif document_type == "csv":
        return extract_csv_text(file_path)

    elif document_type == "txt":
        return extract_txt_text(file_path)

    else:
        print("Unknown Document Type")
        return ""